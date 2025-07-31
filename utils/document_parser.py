# utils/document_parser.py

import os
import PyPDF2 as pypdf2 # For parsing PDF files
from docx import Document # For parsing DOCX files
import pandas as pd # For parsing CSV files
from pptx import Presentation # For parsing PPTX files
# Removed openpyxl as it's not currently used for parsing XLSX

def parse_pdf(file_path: str) -> str:
    """
    Parses text content from a PDF file.

    Args:
        file_path (str): The absolute or relative path to the PDF file.

    Returns:
        str: The extracted text content from the PDF. Returns an empty string
             if the file cannot be opened or parsed.
    """
    text = ""
    try:
        with open(file_path, 'rb') as file:
            # PdfReader handles older and newer PDF formats
            reader = pypdf2.PdfReader(file)
            for page_num in range(len(reader.pages)):
                # extract_text() can return None if a page has no extractable text
                text += reader.pages[page_num].extract_text() or ""
    except Exception as e:
        print(f"[DocumentParser] Error parsing PDF '{file_path}': {e}")
        # Return empty string to signal failure in parsing
        return "" 
    return text

def parse_docx(file_path: str) -> str:
    """
    Parses text content from a DOCX (Microsoft Word) file.

    Args:
        file_path (str): The absolute or relative path to the DOCX file.

    Returns:
        str: The extracted text content from the DOCX. Returns an empty string
             if the file cannot be opened or parsed.
    """
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        print(f"[DocumentParser] Error parsing DOCX '{file_path}': {e}")
        return ""
    return text

def parse_csv(file_path: str) -> str:
    """
    Parses tabular data from a CSV file into a string representation.
    Uses pandas to read the CSV and convert the DataFrame to a string.

    Args:
        file_path (str): The absolute or relative path to the CSV file.

    Returns:
        str: A string representation of the CSV's content. Returns an empty string
             if the file cannot be opened or parsed.
    """
    try:
        df = pd.read_csv(file_path)
        # Convert DataFrame to a string, including headers and index
        # For RAG, a comprehensive string representation is usually preferred.
        return df.to_string(index=False) # index=False makes it cleaner without pandas index
    except Exception as e:
        print(f"[DocumentParser] Error parsing CSV '{file_path}': {e}")
        return ""

def parse_pptx(file_path: str) -> str:
    """
    Parses text content from a PPTX (Microsoft PowerPoint) file.
    Extracts text from all shapes that contain text across all slides.

    Args:
        file_path (str): The absolute or relative path to the PPTX file.

    Returns:
        str: The extracted text content from the PPTX. Returns an empty string
             if the file cannot be opened or parsed.
    """
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame: # Check for text_frame existence
                    text += shape.text_frame.text + "\n"
                elif hasattr(shape, "table") and shape.table: # Handle tables in PPTX
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text += cell.text_frame.text + "\t"
                        text += "\n"
    except Exception as e:
        print(f"[DocumentParser] Error parsing PPTX '{file_path}': {e}")
        return ""
    return text

def parse_txt(file_path: str) -> str:
    """
    Parses text content from a plain TXT file.

    Args:
        file_path (str): The absolute or relative path to the TXT file.

    Returns:
        str: The full text content of the TXT file. Returns an empty string
             if the file cannot be opened or read.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        print(f"[DocumentParser] Error parsing TXT '{file_path}': {e}")
        return ""

def parse_document(file_path: str, file_type: str) -> str:
    """
    Dispatches to the appropriate parser function based on the specified file type.

    Args:
        file_path (str): The path to the document file.
        file_type (str): The type of the file (e.g., 'pdf', 'docx', 'txt', 'csv', 'pptx').
                         Case-insensitive.

    Returns:
        str: The extracted text content of the document. Returns an empty string
             if the file type is unsupported, or if an error occurs during parsing.
    """
    if not os.path.exists(file_path):
        print(f"[DocumentParser] File not found: '{file_path}'")
        return ""

    # Convert to lowercase to ensure case-insensitivity
    file_type_lower = file_type.lower()

    if file_type_lower == 'pdf':
        return parse_pdf(file_path)
    elif file_type_lower == 'docx':
        return parse_docx(file_path)
    elif file_type_lower == 'txt':
        return parse_txt(file_path)
    elif file_type_lower == 'csv':
        return parse_csv(file_path)
    elif file_type_lower == 'pptx':
        return parse_pptx(file_path)
    # Add more parsers here as needed (e.g., for XLSX)
    # elif file_type_lower == 'xlsx':
    #     return parse_xlsx(file_path)
    else:
        print(f"[DocumentParser] Unsupported file type: '{file_type}' for file '{file_path}'.")
        return ""

if __name__ == '__main__':
    # --- Standalone Test for utils/document_parser.py ---

    print("--- Testing utils/document_parser.py ---")

    # Use a temporary directory for test files
    import shutil
    import tempfile
    TEST_DIR = tempfile.mkdtemp()
    print(f"Using temporary directory for DocumentParser tests: {TEST_DIR}")

    # Create dummy files for testing
    txt_file_path = os.path.join(TEST_DIR, "test.txt")
    csv_file_path = os.path.join(TEST_DIR, "test.csv")
    pdf_file_path = os.path.join(TEST_DIR, "test.pdf") # For testing error handling
    docx_file_path = os.path.join(TEST_DIR, "test.docx") # For testing error handling
    pptx_file_path = os.path.join(TEST_DIR, "test.pptx") # For testing error handling
    
    # 1. Test TXT
    print("\n--- Testing TXT parsing ---")
    with open(txt_file_path, "w", encoding='utf-8') as f:
        f.write("This is a test text file.\nIt has multiple lines.\nLine 3.")
    txt_content = parse_document(txt_file_path, "txt")
    print(f"TXT Content (first 50 chars):\n'{txt_content[:50]}...'")
    assert "This is a test text file." in txt_content
    assert "\nLine 3." in txt_content
    print("TXT parsing test PASSED.")

    # 2. Test CSV
    print("\n--- Testing CSV parsing ---")
    with open(csv_file_path, "w", encoding='utf-8') as f:
        f.write("Name,Age,City\nAlice,30,New York\nBob,24,London")
    csv_content = parse_document(csv_file_path, "csv")
    print(f"CSV Content:\n'{csv_content}'")
    assert "Name Age City" in csv_content # pandas to_string typically has a header line
    assert "Alice" in csv_content
    assert "London" in csv_content
    print("CSV parsing test PASSED.")

    # 3. Test non-existent file
    print("\n--- Testing non-existent file ---")
    non_existent_content = parse_document(os.path.join(TEST_DIR, "nonexistent.txt"), "txt")
    print(f"Non-existent file content: '{non_existent_content}'")
    assert non_existent_content == ""
    print("Non-existent file test PASSED.")

    # 4. Test unsupported file type
    print("\n--- Testing unsupported file type ---")
    unsupported_content = parse_document(txt_file_path, "xyz")
    print(f"Unsupported file type content: '{unsupported_content}'")
    assert unsupported_content == ""
    print("Unsupported file type test PASSED.")

    # 5. Test empty/corrupt files (demonstrates error handling)
    print("\n--- Testing empty/corrupt files for error handling ---")
    # For PDF/DOCX/PPTX, creating truly valid but empty/minimal files programmatically
    # can be complex and beyond the scope of a simple utils test.
    # We'll create empty/invalid files and check that parse functions return empty string.

    # Empty PDF (will likely fail parsing)
    with open(pdf_file_path, "wb") as f:
        f.write(b"") # Empty binary file
    corrupt_pdf_content = parse_document(pdf_file_path, "pdf")
    print(f"Corrupt PDF content: '{corrupt_pdf_content}'")
    assert corrupt_pdf_content == ""
    print("Corrupt PDF test PASSED.")

    # Corrupt DOCX (not a real docx structure)
    with open(docx_file_path, "w", encoding='utf-8') as f:
        f.write("This is just plain text, not a DOCX.")
    corrupt_docx_content = parse_document(docx_file_path, "docx")
    print(f"Corrupt DOCX content: '{corrupt_docx_content}'")
    assert corrupt_docx_content == ""
    print("Corrupt DOCX test PASSED.")

    # Corrupt PPTX (not a real pptx structure)
    with open(pptx_file_path, "w", encoding='utf-8') as f:
        f.write("This is just plain text, not a PPTX.")
    corrupt_pptx_content = parse_document(pptx_file_path, "pptx")
    print(f"Corrupt PPTX content: '{corrupt_pptx_content}'")
    assert corrupt_pptx_content == ""
    print("Corrupt PPTX test PASSED.")
    
    # --- Cleanup ---
    print(f"\n--- Cleaning up temporary directory: {TEST_DIR} ---")
    shutil.rmtree(TEST_DIR)
    print("DocumentParser standalone test complete.")